from pptx import Presentation
from pptx.util import Inches

# Create a new presentation object
prs = Presentation()

# ---------------------------------------------------
# Slide 1: Title Slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "Rust System Programming Vulnerabilities Using Dynamic Debugging"
slide.placeholders[1].text = "Techniques & Tools for Identifying Runtime Issues in Rust\nYour Name / Institution"

# ---------------------------------------------------
# Slide 2: Introduction
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Introduction"
content = (
    "• Why Rust?\n"
    "  - Memory safety guarantees\n"
    "  - Growing use in system-level programming\n\n"
    "• Why Debugging and Vulnerabilities Matter?\n"
    "  - Logical or runtime errors can still occur\n"
    "  - Early detection prevents security and stability issues"
)
slide.shapes.placeholders[1].text = content

# ---------------------------------------------------
# Slide 3: Objective
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Objective"
content = (
    "“Cargo test output to debug logs inside test code.”\n\n"
    "Goal:\n"
    " - Demonstrate vulnerability detection in Rust using dynamic debugging tools\n"
    " - Integrate testing (cargo test) with logging\n\n"
    "Key Point:\n"
    " - Utilize Rust's built-in testing framework alongside debuggers and analyzers"
)
slide.shapes.placeholders[1].text = content

# ---------------------------------------------------
# Slide 4: Tools & Purpose (Overview)
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Tools & Purpose (Overview)"
content = (
    "1. GDB\n"
    "   - Step through code, inspect variables, memory, and call stacks\n\n"
    "2. Instrumentation (LLVM & AddressSanitizer)\n"
    "   - Insert compile-time and runtime checks\n"
    "   - Use AddressSanitizer (ASan) to detect memory errors like buffer overflows, use-after-free, etc.\n\n"
    "3. Miri\n"
    "   - Interprets Rust's MIR to catch undefined behaviors and subtle runtime issues\n\n"
    "4. RUST_BACKTRACE\n"
    "   - Displays detailed stack traces on panic, aiding in pinpointing runtime errors"
)
slide.shapes.placeholders[1].text = content

# ---------------------------------------------------
# Slide 5: In-Depth: GDB & Dynamic Debugging
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "In-Depth: GDB & Dynamic Debugging"
content = (
    "GDB Capabilities:\n"
    " - Set breakpoints to halt execution at specific points\n"
    " - Inspect variables and memory values\n"
    " - Trace function calls using call stacks\n\n"
    "Practical Steps:\n"
    " 1. Compile with debugging symbols: cargo build --debug\n"
    " 2. Run your program in GDB: gdb target/debug/<program>\n"
    " 3. Set breakpoints and step through the execution"
)
slide.shapes.placeholders[1].text = content

# ---------------------------------------------------
# Slide 6: Further Instrumentation: LLVM, AddressSanitizer & Miri
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Further Instrumentation: LLVM, AddressSanitizer & Miri"
content = (
    "LLVM Instrumentation:\n"
    " - Insert sanitizers and additional runtime checks\n"
    " - **AddressSanitizer (ASan):**\n"
    "     • Detects memory errors (buffer overflows, use-after-free, etc.) via runtime instrumentation\n\n"
    "Miri:\n"
    " - Interprets Rust's Mid-level Intermediate Representation (MIR) to detect undefined behavior\n"
    " - Catches issues that might be missed at compile time"
)
slide.shapes.placeholders[1].text = content

# ---------------------------------------------------
# Slide 7: Embedded Debugging: probe-rs / JDB
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Embedded Debugging: probe-rs / JDB"
content = (
    "For embedded Rust development:\n\n"
    "probe-rs:\n"
    " - An open-source tool for debugging embedded systems\n"
    " - Interfaces with hardware probes (e.g., J-Link, ST-Link)\n\n"
    "JDB / Related Tools:\n"
    " - Provide standardized interfaces for debugging on microcontrollers\n"
    " - Enable stepping through code and monitoring memory registers"
)
slide.shapes.placeholders[1].text = content

# ---------------------------------------------------
# Slide 8: Argument & Findings
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Argument & Findings"
content = (
    "“The project successfully identified runtime vulnerabilities in Rust by combining dynamic debugging tools, including GDB, AddressSanitizer, Miri, and LLVM instrumentation.”\n\n"
    "Core Argument:\n"
    " - Multiple debugging tools together improve vulnerability detection\n\n"
    "Key Takeaway:\n"
    " - Comprehensive memory analysis increases reliability in system-level programming"
)
slide.shapes.placeholders[1].text = content

# ---------------------------------------------------
# Slide 9: Conclusion
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Conclusion"
content = (
    "Summary:\n"
    " - Rust's inherent safety does not completely eliminate runtime vulnerabilities\n"
    " - Dynamic debugging using multiple tools enhances code security and stability\n\n"
    "Future Considerations:\n"
    " - Incorporate automated testing and CI/CD pipelines with instrumentation\n"
    " - Explore advanced sanitizers, fuzz testing, and further debugging enhancements\n"
    " - Continue refining the Rust debugging ecosystem"
)
slide.shapes.placeholders[1].text = content

# ---------------------------------------------------
# Optional: References Slide
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "References"
content = (
    "• Rust Official Documentation\n"
    "• GDB Documentation\n"
    "• LLVM and AddressSanitizer Documentation\n"
    "• Miri Documentation"
)
slide.shapes.placeholders[1].text = content

# Save the presentation file
ppt_filename = "Rust_System_Programming_Vulnerabilities_Using_Dynamic_Debugging.pptx"
prs.save(ppt_filename)
print(f"Presentation saved as '{ppt_filename}'")
